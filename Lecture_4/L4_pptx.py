import pptx

pre = pptx.Presentation()  # 创建 PPT 对象

# 第一页
title_slide_layout = pre.slide_layouts[0]  # 主题样式序号
slide = pre.slides.add_slide(title_slide_layout)  # 增加到 PPT 对象中

title = slide.shapes.title
subtitle = slide.placeholders[1] # 占位符
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"


# 第二页
bullet_slide_layout = pre.slide_layouts[1]  # 主题样式序号
slide = pre.slides.add_slide(bullet_slide_layout)  # 增加到 PPT 对象中
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Adding a Bullet Slide"
tf = body_shape.text_frame
tf.text = "Find the bullet slide layout"
p = tf.add_paragraph()
p.text = "Use _TextFrame.text for first bullet"
p.level = 1
p = tf.add_paragraph()
p.text = "Use _TextFrame.add_paragraph() for subsequent bullets"
p.level = 2

from pptx.util import Inches

# 第三页
blank_slide_layout = pre.slide_layouts[6]  # 主题样式序号
slide = pre.slides.add_slide(blank_slide_layout)

left = Inches(1)
top = Inches(1)
img_path = "Data_samples/cat.jpg"
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(6)
height = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, height=height)


pre.save("Output/test.pptx")
